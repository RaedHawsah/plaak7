import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "تطبيق ساموراي للتوصيل ⚔️"
    subtitle.text = "أسرع تطبيق توصيل في المملكة\nإعداد الطالب: ............."
    
    # Formatting title
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 204, 0) # Samurai Yellow

    # 2. Problem & Solution Slide
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "المشكلة والحل"
    tf = content.text_frame
    tf.text = "المشكلة:"
    p = tf.add_paragraph()
    p.text = "الحاجة إلى تطبيق توصيل سريع جداً للمطاعم المحلية، بواجهة بسيطة وسهلة الاستخدام تدعم اللغة العربية بالكامل."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "الحل (ساموراي):"
    p2.level = 0
    p3 = tf.add_paragraph()
    p3.text = "تطبيق صُمم ليكون سريعاً كدقة الساموراي، يركز على الوجبات المفضلة مع ميزات تفاعلية حديثة وسرعة استجابة عالية (بدون انتظار تحميل طويل)."
    p3.level = 1

    # 3. Features Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "المميزات الرئيسية"
    tf = content.text_frame
    tf.text = "دعم كامل للغة العربية (RTL) وتصميم احترافي."
    p = tf.add_paragraph()
    p.text = "نظام سلة مشتريات (Cart) متكامل لتحديث الكميات وإتمام الطلب."
    p = tf.add_paragraph()
    p.text = "لوحة تحكم للمطاعم (محاكاة): نظام فلترة وترتيب الوجبات بناءً على السعر."
    p = tf.add_paragraph()
    p.text = "أكواد خصم (Promo Codes) تفاعلية يمكن نسخها واستخدامها."
    p = tf.add_paragraph()
    p.text = "نظام تسجيل دخول وهمي وحفظ السجل للمستخدم محلياً."

    # 4. Technologies Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "التقنيات والأدوات المستخدمة"
    tf = content.text_frame
    tf.text = "تم بناء المشروع بدون إطارات عمل لإثبات كفاءة التعامل مع الأساسيات:"
    p = tf.add_paragraph()
    p.text = "HTML5: لبناء الصفحات والهيكلة الدلالية (Semantic HTML)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "CSS3: للتصميم المرن (Flexbox/Grid)، المتغيرات (Variables)، والحركات التفاعلية."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "JavaScript (ES6): للتفاعلية، تحديث واجهة المستخدم (DOM)، وحساب السلة."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "LocalStorage API: لمحاكاة قاعدة البيانات وتخزين الطلبات والمستخدمين."
    p.level = 1

    # 5. Conclusion Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "شاشة العرض (Demo)"
    tf = content.text_frame
    tf.text = "التطبيق جاهز للاستعراض المباشر (Live Demo)."
    p = tf.add_paragraph()
    p.text = "شكراً لحسن استماعكم."
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "مستعد للإجابة على استفساراتكم."
    p2.level = 0

    prs.save('Presentation.pptx')
    print("Presentation.pptx generated successfully.")

if __name__ == '__main__':
    create_ppt()
